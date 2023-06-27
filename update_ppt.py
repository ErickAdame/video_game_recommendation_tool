from pptx import Presentation
import openpyxl

# Set the file paths
excel_file_path = "/Users/erick/Desktop/city_high_and Lows.csv"
powerpoint_file_path = "/Users/erick/Desktop/Weather Update.pptx"

# Set the cell reference in Excel to retrieve the data from
excel_cell = "Sheet1!C2"

# Set the slide index and text box index of the PowerPoint slide to update
slide_index = 7  # Change this to the desired slide index (0-based)
textbox_index = 9  # Change this to the desired text box index (0-based)

# Load the Excel workbook and retrieve the cell value
excel_workbook = openpyxl.load_workbook(excel_file_path)
excel_sheet = excel_workbook[excel_cell.split('!')[0]]
cell_value = excel_sheet[excel_cell.split('!')[1]].value

# Load the PowerPoint presentation
presentation = Presentation(powerpoint_file_path)

# Access the desired slide and text box to update the content
slide = presentation.slides[slide_index]
textbox = slide.shapes[textbox_index].text_frame

# Clear existing text and set the new value
textbox.clear()
textbox.text = str(cell_value)

# Save the updated presentation
presentation.save("updated_powerpoint.pptx")
